"""Embed mono + poly WAV files as click-to-play audio on Slide 5 of the P2 deck.

Slide 5 layout (extracted via python-pptx):
  Mono panel image: left=9.18 in, top=2.61 in, w=4.89 x h=3.06 in
  Poly panel image: left=14.27 in, top=2.61 in, w=4.89 x h=3.06 in

We place a small speaker icon at the top-left of each panel, slightly inset
so it's visible without covering the spectrogram. Each icon is 0.45 x 0.45 in.

PowerPoint renders these as click-to-play audio (no auto-play). Verify after
opening in PowerPoint and pressing F5 -> click each icon.

Run:
    python deliverables/p2/embed_audio.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent
DECK = ROOT / "P2-Presentation.pptx"
MONO_WAV = ROOT / "audio_samples" / "mono_example.wav"
POLY_WAV = ROOT / "audio_samples" / "poly_example.wav"


def main() -> None:
    for p in (DECK, MONO_WAV, POLY_WAV):
        if not p.exists():
            raise SystemExit(f"missing file: {p}")

    prs = Presentation(str(DECK))
    slide = prs.slides[4]  # 0-indexed slide 5

    # Remove any prior audio shapes (re-runnable).
    # Audio shapes have shape_type = MEDIA (16) in PPTX.
    for shape in list(slide.shapes):
        if shape.shape_type == 16:
            sp_el = shape._element
            sp_el.getparent().remove(sp_el)
            print(f"removed existing audio shape {shape.name!r}")

    # Mono — top-left corner of mono panel (panel at 9.18, 2.61)
    mono_icon = slide.shapes.add_movie(
        str(MONO_WAV),
        left=Inches(9.30),
        top=Inches(2.73),
        width=Inches(0.45),
        height=Inches(0.45),
        mime_type="audio/wav",
    )
    print(f"added MONO icon at (9.30, 2.73) -> {mono_icon.name}")

    # Poly — top-left corner of poly panel (panel at 14.27, 2.61)
    poly_icon = slide.shapes.add_movie(
        str(POLY_WAV),
        left=Inches(14.39),
        top=Inches(2.73),
        width=Inches(0.45),
        height=Inches(0.45),
        mime_type="audio/wav",
    )
    print(f"added POLY icon at (14.39, 2.73) -> {poly_icon.name}")

    prs.save(str(DECK))
    print(f"\nsaved {DECK.name} with embedded audio.")
    print("Open in PowerPoint, press F5, click each icon to verify click-to-play.")


if __name__ == "__main__":
    main()
